#!/usr/bin/env python3
from __future__ import annotations

import csv
import json
import os
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
RUN_DIR = (
    ROOT / "experiments" / "runs" / "real_api_50_full_after_fix"
    if (ROOT / "experiments" / "runs" / "real_api_50_full_after_fix").exists()
    else (
        ROOT / "experiments" / "runs" / "submission_batch_verify"
        if (ROOT / "experiments" / "runs" / "submission_batch_verify").exists()
        else ROOT / "experiments" / "runs" / "batch_verify"
    )
)
ASSET_DIR = ROOT / "deliverables" / "assets"
GALLERY_DIR = ROOT / "deliverables" / "blender_gallery"
OUT_DIR = ROOT / "deliverables" / "slides"
PPTX_PATH = OUT_DIR / "SceneTest_presentation.pptx"
MANIFEST_PATH = OUT_DIR / "presentation_manifest.json"

CODE_URL = os.environ.get("SCENETEST_CODE_URL", "https://github.com/16yunH/SceneTest-PJ3")
AUTHOR = f"{os.environ.get('SCENETEST_AUTHOR_CN', '洪运')}  {os.environ.get('SCENETEST_STUDENT_ID', '23300240019')}"
FONT = "PingFang SC"
MONO_FONT = "Menlo"

W, H = 13.333, 7.5
COLORS = {
    "bg": "F6F8FB",
    "ink": "18212B",
    "muted": "5B6675",
    "panel": "FFFFFF",
    "line": "D5DBE5",
    "blue": "2563A7",
    "teal": "237A74",
    "amber": "B8741A",
    "red": "B84646",
    "green": "2E7D4F",
    "dark": "1F2A37",
    "soft_blue": "EAF2FB",
    "soft_teal": "E7F4F1",
    "soft_amber": "FFF0DB",
    "soft_red": "FCEBEB",
    "code": "111827",
}


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.strip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(0.7)
    else:
        shp.line.fill.background()
    return shp


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=16,
    color=None,
    bold=False,
    align="left",
    font=FONT,
    valign="top",
    leading=1.08,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(valign, MSO_ANCHOR.TOP)
    alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = alignment
        p.line_spacing = leading
        p.font.name = font
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color or COLORS["ink"])
    return box


def add_bullets(slide, items, x, y, w, h, size=15, color=None):
    text = "\n".join(f"• {item}" for item in items)
    return add_text(slide, text, x, y, w, h, size=size, color=color or COLORS["muted"], leading=1.14)


def base_slide(slide, eyebrow, title, subtitle=""):
    add_rect(slide, 0, 0, W, H, COLORS["bg"])
    add_rect(slide, 0, 0, W, 0.08, COLORS["dark"])
    add_text(slide, eyebrow, 0.72, 0.38, 5.8, 0.28, 12.5, COLORS["blue"], True)
    add_text(slide, title, 0.72, 0.72, 10.5, 0.48, 27, COLORS["ink"], True)
    if subtitle:
        add_text(slide, subtitle, 0.72, 1.24, 10.7, 0.38, 14.5, COLORS["muted"])
    add_text(slide, "SceneTest | 计算机图形学 PJ3", 10.1, 6.92, 2.45, 0.22, 9.5, COLORS["muted"], align="right")


def panel(slide, x, y, w, h, title, body, fill=None, accent=None):
    add_rect(slide, x, y, w, h, fill or COLORS["panel"], COLORS["line"], radius=True)
    if accent:
        add_rect(slide, x, y, 0.08, h, accent)
    if h <= 0.9:
        add_text(slide, title, x + 0.18, y + 0.12, w - 0.36, 0.24, 13.6, accent or COLORS["ink"], True)
        add_text(slide, body, x + 0.18, y + 0.42, w - 0.36, max(0.2, h - 0.5), 10.7, COLORS["muted"], leading=1.05)
    else:
        add_text(slide, title, x + 0.18, y + 0.16, w - 0.36, 0.32, 15.2, accent or COLORS["ink"], True)
        add_text(slide, body, x + 0.18, y + 0.56, w - 0.36, h - 0.68, 12.8, COLORS["muted"], leading=1.15)


def metric_card(slide, x, y, value, label, color):
    add_rect(slide, x, y, 2.35, 1.0, COLORS["panel"], COLORS["line"], radius=True)
    add_rect(slide, x, y, 0.08, 1.0, color)
    add_text(slide, value, x + 0.22, y + 0.12, 1.9, 0.38, 27, color, True)
    add_text(slide, label, x + 0.22, y + 0.62, 1.85, 0.24, 11.2, COLORS["muted"])


def image_fit(slide, path, x, y, w, h):
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    return slide.shapes.add_picture(str(path), Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2), Inches(pw), Inches(ph))


def image_panel(slide, x, y, w, h, path, label):
    add_rect(slide, x, y, w, h, COLORS["panel"], COLORS["line"], radius=True)
    image_fit(slide, path, x + 0.16, y + 0.16, w - 0.32, h - 0.58)
    add_text(slide, label, x + 0.18, y + h - 0.36, w - 0.36, 0.22, 12.3, COLORS["ink"], True, "center")


def gif_panel(slide, x, y, w, h, path, label):
    image_panel(slide, x, y, w, h, path, f"{label}  |  GIF")


def code_block(slide, x, y, w, h, title, code, size=10.5):
    add_rect(slide, x, y, w, h, COLORS["code"], radius=True)
    add_text(slide, title, x + 0.18, y + 0.15, w - 0.36, 0.22, 10.5, "9CA3AF", True)
    add_text(slide, code, x + 0.18, y + 0.48, w - 0.36, h - 0.62, size, "E5E7EB", font=MONO_FONT, leading=1.02)


def flow(slide, labels, x, y, box_w=1.55, box_h=0.72):
    cur = x
    for i, label in enumerate(labels):
        fill = COLORS["soft_teal"] if i % 2 else COLORS["panel"]
        add_rect(slide, cur, y, box_w, box_h, fill, COLORS["line"], radius=True)
        add_text(slide, label, cur + 0.08, y + 0.22, box_w - 0.16, 0.24, 12.2, COLORS["ink"], True, "center")
        if i < len(labels) - 1:
            add_text(slide, "→", cur + box_w + 0.03, y + 0.18, 0.28, 0.28, 17, COLORS["muted"], True, "center")
        cur += box_w + 0.34


def test_row(slide, x, y, w, title, check, source, fill):
    add_rect(slide, x, y, w, 0.48, fill, COLORS["line"], radius=True)
    add_rect(slide, x, y, 0.07, 0.48, COLORS["dark"])
    add_text(slide, title, x + 0.18, y + 0.12, 1.05, 0.2, 13.2, COLORS["ink"], True)
    add_text(slide, check, x + 1.42, y + 0.14, 2.1, 0.18, 10.8, COLORS["muted"], font=MONO_FONT)
    add_text(slide, source, x + w + 0.34, y + 0.14, 2.35, 0.18, 11.2, COLORS["muted"])


def read_rows():
    with (RUN_DIR / "results.csv").open(newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def read_llm_audit_sample():
    path = RUN_DIR / "_llm_audit_sample" / "summary.json"
    if not path.exists():
        return {
            "sample_size": 0,
            "pass_count": 0,
            "avg_overall_score": 0.0,
            "avg_elapsed_sec": 0.0,
            "estimated_full_50_minutes": 0.0,
            "rows": [],
        }
    return json.loads(path.read_text(encoding="utf-8"))


def read_timing_comparison():
    path = RUN_DIR / "_llm_audit_sample" / "timing_comparison.json"
    if not path.exists():
        return {
            "no_llm_audit_avg_sec": 0.0,
            "llm_audit_overhead_avg_sec": 0.0,
            "with_llm_audit_avg_sec": 0.0,
            "estimated_50_no_audit_minutes": 0.0,
            "estimated_50_with_audit_minutes": 0.0,
        }
    return json.loads(path.read_text(encoding="utf-8"))


def metric(rows, method, name):
    for row in rows:
        if row["method"] == method and row["metric"] == name:
            return float(row["rate"]), int(row["passed"]), int(row["total"])
    return 0.0, 0, 0


def pct(v):
    return f"{round(v * 100)}%"


def bar_chart(slide, x, y, w, h, title, rows, metric_name):
    add_rect(slide, x, y, w, h, COLORS["panel"], COLORS["line"], radius=True)
    add_text(slide, title, x + 0.2, y + 0.18, w - 0.4, 0.28, 14.2, COLORS["ink"], True)
    methods = [
        ("单次生成", "single_pass", COLORS["red"]),
        ("契约生成", "contract_only", COLORS["amber"]),
        ("SceneTest", "scenetest", COLORS["green"]),
    ]
    for i, (label, key, color) in enumerate(methods):
        value, _, _ = metric(rows, key, metric_name)
        yy = y + 0.72 + i * 0.58
        add_text(slide, label, x + 0.22, yy, 0.82, 0.22, 10.8, COLORS["muted"])
        add_rect(slide, x + 1.12, yy + 0.04, w - 1.66, 0.15, "E7ECF3", radius=True)
        add_rect(slide, x + 1.12, yy + 0.04, max(0.05, (w - 1.66) * value), 0.15, color, radius=True)
        add_text(slide, pct(value), x + w - 0.48, yy - 0.01, 0.38, 0.22, 10.8, COLORS["ink"], True, "right")


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    rows = read_rows()
    audit_sample = read_llm_audit_sample()
    timing = read_timing_comparison()
    final_rate, final_passed, final_total = metric(rows, "scenetest", "Overall Contract Pass Rate")
    single_rate, _, _ = metric(rows, "single_pass", "Overall Contract Pass Rate")
    contract_rate, _, _ = metric(rows, "contract_only", "Overall Contract Pass Rate")
    relation_rate, relation_passed, relation_total = metric(rows, "scenetest", "Spatial Relation Accuracy")
    object_rate, object_passed, object_total = metric(rows, "scenetest", "Object Completeness")
    visible_rate, visible_passed, visible_total = metric(rows, "scenetest", "Visibility Pass Rate")
    contract = json.loads((RUN_DIR / "desk_cozy" / "contract.json").read_text(encoding="utf-8"))
    repair = json.loads((RUN_DIR / "desk_cozy" / "scenetest" / "repair_history.json").read_text(encoding="utf-8"))

    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, W, H, COLORS["dark"])
    add_rect(slide, 0, 0, 4.55, H, "EEF3F8")
    add_text(slide, "SceneTest", 0.66, 0.72, 3.55, 0.64, 43, COLORS["ink"], True)
    add_text(slide, "场景契约 + 图形单元测试\n驱动的 Agentic 文本到 3D 场景生成", 0.70, 1.58, 3.55, 0.92, 18.8, COLORS["dark"], True)
    add_text(slide, "计算机图形学 Project 3", 0.72, 5.85, 3.4, 0.24, 13, COLORS["dark"])
    add_text(slide, AUTHOR, 0.72, 6.18, 3.4, 0.24, 13.5, COLORS["dark"], True)
    add_text(slide, CODE_URL, 0.72, 6.48, 3.5, 0.22, 9.3, COLORS["blue"])
    image_fit(slide, ASSET_DIR / "deepseek_blender.png", 5.15, 1.05, 6.75, 4.25)
    metric_card(slide, 5.25, 5.72, f"{final_passed}/{final_total}", "修复后测试全部通过", COLORS["green"])
    metric_card(slide, 8.05, 5.72, f"+{(final_rate - single_rate) * 100:.1f}pt", "总体通过率提升", COLORS["blue"])

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "01 研究动机", "LLM 生成 3D 场景时，错误是可测试的", "场景看起来合理，并不代表它满足 prompt 中的显式要求。")
    panel(slide, 0.82, 2.0, 2.55, 1.55, "对象缺失", "prompt 要求植物或杯子，但生成代码没有创建对应物体。", COLORS["panel"], COLORS["red"])
    panel(slide, 3.65, 2.0, 2.55, 1.55, "关系错误", "杯子应在电脑右侧，却出现在左侧、重叠或悬空。", COLORS["soft_amber"], COLORS["amber"])
    panel(slide, 6.48, 2.0, 2.55, 1.55, "不可见", "物体存在于场景状态中，但相机画面里看不到。", COLORS["soft_blue"], COLORS["blue"])
    panel(slide, 9.31, 2.0, 2.55, 1.55, "修复不稳", "整段重写可能修好一个问题，同时引入新错误。", COLORS["soft_teal"], COLORS["teal"])
    add_text(slide, "核心观察：既然软件可以用单元测试约束行为，3D 场景生成也可以把 prompt 要求编译成可执行的图形测试。", 1.18, 4.65, 10.9, 0.68, 21, COLORS["ink"], True, "center")

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "02 核心想法", "把 prompt 对齐变成“可测试、可修复”的图形编程流程", "SceneTest 在生成和修复之间加入结构化契约与测试层。")
    flow(slide, ["用户 Prompt", "Scene Contract", "场景代码", "图形单元测试", "失败报告", "局部修复"], 0.68, 2.45, 1.55)
    panel(slide, 1.12, 4.35, 4.6, 1.25, "契约层", "对象、空间关系、材质、光照、相机可见目标都写入同一份 JSON。", COLORS["soft_teal"], COLORS["teal"])
    panel(slide, 7.02, 4.35, 4.6, 1.25, "修复层", "每个失败测试映射到明确动作：创建物体、移动位置、设置光照、调整相机。", COLORS["soft_blue"], COLORS["blue"])

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "03 系统架构", "三个 Agent 加一个测试闭环，再加按需语义审查", "DeepSeek 用于真实 Contract Agent 和可选 LLM semantic audit；离线 parser 保证仓库可复现。")
    flow(slide, ["Contract Agent", "Code Agent", "SceneBuilder", "Test Compiler", "Test Runner", "Repair Agent"], 0.75, 2.24, 1.48)
    add_text(slide, "↘  可选：LLM Semantic Audit 读取 prompt / contract / scene.json / test_results", 2.0, 3.15, 8.6, 0.28, 13.2, COLORS["blue"], True, "center")
    panel(slide, 0.92, 4.0, 3.25, 1.45, "共享语义", "所有模块读取同一份 Scene Contract，避免各自重新猜 prompt。", COLORS["panel"], COLORS["blue"])
    panel(slide, 4.92, 4.0, 3.25, 1.45, "受控生成", "Code Agent 只生成 helper API 调用，而不是自由写任意 bpy 代码。", COLORS["soft_teal"], COLORS["teal"])
    panel(slide, 8.92, 4.0, 3.25, 1.45, "闭环验证", "测试失败后只做局部修复，再重新测试，最多 3 轮。", COLORS["soft_amber"], COLORS["amber"])

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "04 场景契约", "自然语言被编译成可度量的 3D 需求", "契约同时服务于代码生成、测试编译和修复规则。")
    contract_code = "\n".join(
        [
            "{",
            '  "objects": ["desk", "laptop", "lamp",',
            '              "coffee_cup", "plant"],',
            '  "relations": [',
            '    "coffee_cup right_of laptop",',
            '    "lamp left_of laptop",',
            '    "plant behind laptop",',
            '    "laptop on desk"',
            "  ],",
            '  "style": {"lighting": "warm"},',
            '  "camera": {"visible": "all required objects"}',
            "}",
        ]
    )
    code_block(slide, 0.82, 1.92, 6.45, 4.2, "desk_cozy contract.json", contract_code, 10.4)
    panel(slide, 7.8, 2.2, 3.7, 0.78, "Objects", "稳定 id + 类型 + 材质/颜色", COLORS["panel"], COLORS["blue"])
    panel(slide, 7.8, 3.25, 3.7, 0.78, "Relations", "on / left_of / right_of / behind", COLORS["soft_teal"], COLORS["teal"])
    panel(slide, 7.8, 4.3, 3.7, 0.78, "Style & Camera", "warm lighting + visible objects", COLORS["soft_amber"], COLORS["amber"])

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "05 图形单元测试", "用几何、外观、光照和相机约束场景", "测试后端是可复现的 SceneBuilder；Blender 用于最终渲染展示。")
    tests = [
        ("对象测试", "assert_exists(laptop)", "scene registry"),
        ("空间关系", "assert_right_of(cup, laptop)", "bbox center"),
        ("支撑关系", "assert_on(laptop, desk)", "z + xy overlap"),
        ("可见性", "assert_visible(cup)", "camera frame"),
        ("风格测试", "assert_lighting(warm)", "light metadata"),
    ]
    for i, (name, check, source) in enumerate(tests):
        yy = 1.95 + i * 0.68
        test_row(slide, 0.85, yy, 4.35, name, check, source, COLORS["soft_teal"] if i % 2 else COLORS["panel"])
    code_block(slide, 8.4, 2.08, 3.25, 2.25, "bbox relation", "left_of:\nsubject.cx < target.cx - margin\n\non:\nabs(subject.min_z - target.max_z) <= 0.08\nxy_overlap >= 0.2", 11.2)

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "06 失败驱动修复", "失败测试直接变成局部修复动作", "不是整段重写场景，而是定位失败项后修改必要状态。")
    failure = json.dumps(repair[0]["failures"][0], ensure_ascii=False, indent=2)
    code_block(slide, 0.78, 1.92, 5.75, 3.3, "failure report", failure, 8.2)
    flow(slide, ["失败测试", "repair_hint", "局部动作", "重新测试"], 1.16, 5.68, 1.8)
    panel(slide, 7.05, 2.18, 4.05, 0.78, "缺失对象", "从 contract 中读取类型并创建 primitive。", COLORS["panel"], COLORS["blue"])
    panel(slide, 7.05, 3.35, 4.05, 0.78, "关系错误", "根据 bounding box 重新设置 x / y / z。", COLORS["soft_teal"], COLORS["teal"])
    panel(slide, 7.05, 4.52, 4.05, 0.78, "不可见", "自动 frame 所有 required objects。", COLORS["soft_amber"], COLORS["amber"])

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "07 案例展示", "同一份契约驱动测试和修复", "Cozy desk 场景：电脑、台灯、杯子、植物、暖光和可见性要求。")
    image_panel(slide, 0.82, 1.9, 5.15, 3.85, ASSET_DIR / "desk_before.png", "修复前：contract-only")
    image_panel(slide, 7.15, 1.9, 5.15, 3.85, ASSET_DIR / "desk_after.png", "修复后：SceneTest")
    add_text(slide, "修复动作：移动 coffee_cup 到 laptop 右侧；设置 warm lighting；重新约束空间关系和相机框选。", 1.45, 6.08, 10.5, 0.28, 14.5, COLORS["ink"], True, "center")

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "08 定量结果", "50 个 prompt、850 个测试，修复后全部通过", "对比 single-pass、contract-only 和完整 SceneTest。")
    bar_chart(slide, 0.68, 1.9, 3.72, 3.0, "总体契约通过率", rows, "Overall Contract Pass Rate")
    bar_chart(slide, 4.82, 1.9, 3.72, 3.0, "空间关系准确率", rows, "Spatial Relation Accuracy")
    bar_chart(slide, 8.96, 1.9, 3.72, 3.0, "可见性通过率", rows, "Visibility Pass Rate")
    metric_card(slide, 2.05, 5.55, str(object_total), "object tests", COLORS["blue"])
    metric_card(slide, 5.48, 5.55, str(relation_total), "relation tests", COLORS["amber"])
    metric_card(slide, 8.91, 5.55, str(visible_total), "visibility tests", COLORS["teal"])

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "09 现场演示", "助教可以现场输入 prompt，直接看到完整闭环", "网页入口调用同一套 CLI pipeline，并读取本地 DeepSeek 配置。")
    code_block(slide, 0.78, 1.85, 5.55, 2.45, "启动命令", "cd SceneTest-PJ3\n.venv/bin/python main.py demo --port 7860\n\nOpen: http://127.0.0.1:7860\nBackend: DeepSeek from local config", 10.6)
    image_panel(slide, 7.05, 1.8, 4.85, 3.45, ASSET_DIR / "deepseek_blender.png", "DeepSeek + SceneTest + Blender")
    panel(slide, 0.92, 4.82, 5.1, 1.12, "网页展示内容", "contract JSON、测试通过率、失败列表、repair history、三组 render 与 Blender 最终图。", COLORS["soft_teal"], COLORS["teal"])
    add_rect(slide, 7.05, 5.58, 4.85, 0.78, COLORS["panel"], COLORS["line"], radius=True)
    add_rect(slide, 7.05, 5.58, 0.08, 0.78, COLORS["blue"])
    add_text(slide, "公开代码", 7.25, 5.72, 1.0, 0.22, 13.6, COLORS["blue"], True)
    add_text(slide, CODE_URL, 7.25, 6.05, 4.35, 0.18, 9.2, COLORS["muted"])

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "10 LLM 语义审查", "按需启用的高信号质检层，而不是主闭环瓶颈", "它不看图，只读取 prompt、contract、scene.json、test results 和 repair history。")
    sample_size = int(audit_sample.get("sample_size", 0) or 0)
    pass_count = int(audit_sample.get("pass_count", 0) or 0)
    avg_score = float(audit_sample.get("avg_overall_score", 0) or 0)
    no_audit = float(timing.get("no_llm_audit_avg_sec", 0) or 0)
    with_audit = float(timing.get("with_llm_audit_avg_sec", 0) or 0)
    audit_overhead = float(timing.get("llm_audit_overhead_avg_sec", 0) or 0)
    no_audit_50 = float(timing.get("estimated_50_no_audit_minutes", 0) or 0)
    with_audit_50 = float(timing.get("estimated_50_with_audit_minutes", 0) or 0)
    metric_card(slide, 0.88, 1.88, f"{no_audit:.0f}s", "不加 audit", COLORS["blue"])
    metric_card(slide, 3.72, 1.88, f"{with_audit:.0f}s", "加 audit", COLORS["teal"])
    metric_card(slide, 6.56, 1.88, f"+{audit_overhead:.0f}s", "audit 增量", COLORS["amber"])
    metric_card(slide, 9.40, 1.88, f"{pass_count}/{sample_size}", "抽样语义通过", COLORS["green"])
    panel(slide, 0.92, 3.45, 3.55, 1.35, "真实计时口径", f"2 个 DeepSeek demo 样本：50 场景估算约 {no_audit_50:.1f}min vs {with_audit_50:.1f}min。", COLORS["panel"], COLORS["blue"])
    panel(slide, 4.88, 3.45, 3.55, 1.35, "宽松审查口径", "承认当前不是精细建模系统；primitive proxy 细节不足只作为建议。", COLORS["soft_teal"], COLORS["teal"])
    panel(slide, 8.84, 3.45, 3.55, 1.35, "主要检查什么", "缺失对象、错误语义类型、支撑关系、空间关系和相对尺度。", COLORS["soft_amber"], COLORS["amber"])
    add_text(slide, f"10 场景抽样平均分 {avg_score:.3f}；{pass_count}/{sample_size} 通过。结论：LLM semantic audit 适合作为按需语义 QA 和展示加分项，而不是每次生成都必须等待的核心路径。", 1.35, 5.52, 10.7, 0.7, 15.5, COLORS["ink"], True, "center")

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "11 Blender GIF 展示", "从纯文本生成可旋转查看的 3D 场景", "PPT 中嵌入 turntable GIF；静态 PDF 中对应为首帧。")
    gif_panel(slide, 0.72, 1.78, 3.86, 3.2, GALLERY_DIR / "drone_repair.gif", "drone_repair")
    gif_panel(slide, 4.74, 1.78, 3.86, 3.2, GALLERY_DIR / "robotics_lab.gif", "robotics_lab")
    gif_panel(slide, 8.76, 1.78, 3.86, 3.2, GALLERY_DIR / "desk_cozy.gif", "desk_cozy")
    add_text(slide, "这些 GIF 来自同一批 50-prompt 实验的最终 SceneTest 场景，不是手工建模素材。", 1.15, 5.58, 11.0, 0.34, 16, COLORS["ink"], True, "center")

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "12 参考项目对比", "我们的创新点在“可执行契约 + 局部修复 + 按需语义审查”", "不是只做一个 prompt-to-Blender demo。")
    panel(slide, 0.82, 1.78, 2.55, 2.55, "直接 LLM→Blender", "优点：快、灵活。\n问题：缺少可复查的失败谓词，修复常依赖整段重写。", COLORS["panel"], COLORS["muted"])
    panel(slide, 3.65, 1.78, 2.55, 2.55, "LL3M / 代码生成式 3D", "侧重语言模型生成 3D 程序。\n我们补充：把 prompt 要求编译成测试。", COLORS["soft_blue"], COLORS["blue"])
    panel(slide, 6.48, 1.78, 2.55, 2.55, "SAGE / Agentic 3D", "强调 agent 级场景生成。\n我们补充：轻量、可复现、课程项目可跑通。", COLORS["soft_teal"], COLORS["teal"])
    panel(slide, 9.31, 1.78, 2.55, 2.55, "SceneTest", "契约驱动测试、失败定位、局部修复、文本语义抽样审查。", COLORS["soft_amber"], COLORS["amber"])
    add_text(slide, "核心差异：我们不把最终准确率建立在主观截图判断上，而是先用代码测试保证硬约束，再用 LLM audit 抽查软语义质量。", 1.12, 5.25, 11.1, 0.58, 18, COLORS["ink"], True, "center")

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "13 可靠性与限制", "当前版本强调可复现闭环，而不是追求最高视觉真实感", "这些限制已写入报告，避免过度声称。")
    panel(slide, 0.92, 2.05, 3.35, 2.25, "为什么结果可信", "测试来自 contract，failure report 可复查；benchmark 可一键复跑。", COLORS["panel"], COLORS["blue"])
    panel(slide, 5.0, 2.05, 3.35, 2.25, "实现取舍", "主测试后端是 SceneBuilder；Blender 用于最终渲染；LLM audit 按需启用。", COLORS["soft_teal"], COLORS["teal"])
    panel(slide, 9.08, 2.05, 3.35, 2.25, "未来增强", "加入 Blender segmentation mask、遮挡判断、更开放的物体库和视觉审查。", COLORS["soft_amber"], COLORS["amber"])

    slide = prs.slides.add_slide(blank)
    base_slide(slide, "14 总结", "SceneTest 让文本到 3D 场景生成变得可检查、可测试、可修复", "项目满足 Agent + 图形学主题，并提供实现、实验、公开代码和现场演示。")
    panel(slide, 0.72, 2.05, 2.85, 1.72, "贡献 1：Scene Contract", "把自然语言需求转成结构化约束。", COLORS["panel"], COLORS["blue"])
    panel(slide, 3.86, 2.05, 2.85, 1.72, "贡献 2：Graphics Unit Tests", "验证对象、关系、可见性、材质和光照。", COLORS["soft_teal"], COLORS["teal"])
    panel(slide, 7.0, 2.05, 2.85, 1.72, "贡献 3：Repair Loop", "失败测试触发局部修复并重新测试。", COLORS["soft_amber"], COLORS["amber"])
    panel(slide, 10.14, 2.05, 2.48, 1.72, "贡献 4：LLM Audit", "按需抽查软语义质量。", COLORS["soft_blue"], COLORS["blue"])
    add_text(slide, CODE_URL, 1.65, 5.2, 10.0, 0.3, 15.5, COLORS["blue"], True, "center")
    add_text(slide, AUTHOR, 1.65, 5.65, 10.0, 0.28, 14, COLORS["muted"], True, "center")

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)
    MANIFEST_PATH.write_text(
        json.dumps(
            {
                "pptx": str(PPTX_PATH.relative_to(ROOT)),
                "slideCount": len(prs.slides),
                "generatedBy": "scripts/build_presentation.py",
                "language": "zh-CN",
                "author": AUTHOR,
                "codeUrl": CODE_URL,
            },
            ensure_ascii=False,
            indent=2,
        )
        + "\n",
        encoding="utf-8",
    )
    print(json.dumps({"pptx": str(PPTX_PATH), "slideCount": len(prs.slides)}, ensure_ascii=False))


if __name__ == "__main__":
    build()
